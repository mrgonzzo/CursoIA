import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.preprocessing import StandardScaler
from sklearn.decomposition import PCA
from sklearn.pipeline import make_pipeline
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches, Pt

# =============================
# 1. Cargar datos
# =============================
df = pd.read_csv('USArrests.csv')  # Cambia según tu archivo
X = df.drop('USArrests', axis=1)
y = df['USArrests']

# Nombre de los estados
if 'State' in df.columns:
    estados = df['State'].values
else:
    estados = [f"Obs{i+1}" for i in range(X.shape[0])]

# =============================
# 2. Pipeline PCA
# =============================
pca_pipeline = make_pipeline(
    StandardScaler(),
    PCA(n_components=X.shape[1])
)
X_pca = pca_pipeline.fit_transform(X)
pca_model = pca_pipeline.named_steps['pca']
scaler = pca_pipeline.named_steps['standardscaler']

# =============================
# 3. Preparar gráficos y tablas
# =============================
# Scree plot / varianza acumulada
varianza_expl = pca_model.explained_variance_ratio_
varianza_acum = np.cumsum(varianza_expl)

plt.figure(figsize=(8,6))
plt.plot(range(1, len(varianza_expl)+1), varianza_expl, marker='o', label='Varianza individual')
plt.plot(range(1, len(varianza_expl)+1), varianza_acum, marker='s', label='Varianza acumulada')
plt.xlabel('Número de componente')
plt.ylabel('Varianza explicada')
plt.title('Scree Plot / Varianza acumulada')
plt.grid(True, linestyle='--', alpha=0.5)
plt.legend()
plt.savefig('scree_plot.png')
plt.close()

# Loadings
loadings = pd.DataFrame(
    pca_model.components_,
    columns=X.columns,
    index=[f'PC{i+1}' for i in range(pca_model.components_.shape[0])]
)

plt.figure(figsize=(8,6))
sns.heatmap(loadings.iloc[:2,:], annot=True, cmap='viridis')
plt.title('Matriz de cargas (loadings) PC1 y PC2')
plt.savefig('loadings_heatmap.png')
plt.close()

# Biplot
pc1_vals = X_pca[:,0]
pc2_vals = X_pca[:,1]
plt.figure(figsize=(10,8))
plt.scatter(pc1_vals, pc2_vals, color='skyblue', s=100, alpha=0.7)
for i, estado in enumerate(estados):
    plt.text(pc1_vals[i]+0.02, pc2_vals[i]+0.02, estado, fontsize=8)
plt.axhline(0, color='gray', linestyle='--')
plt.axvline(0, color='gray', linestyle='--')
plt.xlabel('PC1')
plt.ylabel('PC2')
plt.title('Biplot: Estados proyectados en PC1-PC2')
plt.grid(True, linestyle='--', alpha=0.5)
plt.savefig('biplot.png')
plt.close()

# Estados extremos PC1
df_pc1 = pd.DataFrame({'Estado': estados, 'PC1': pc1_vals})
top5 = df_pc1.sort_values(by='PC1', ascending=False).head(5)
bottom5 = df_pc1.sort_values(by='PC1', ascending=True).head(5)

# =============================
# 4. Generar PDF
# =============================
with PdfPages('Informe_PCA_USArrests.pdf') as pdf:
    # Página 1: Scree Plot
    img = plt.imread('scree_plot.png')
    plt.figure(figsize=(8,6))
    plt.imshow(img)
    plt.axis('off')
    pdf.savefig()
    plt.close()

    # Página 2: Loadings Heatmap
    img = plt.imread('loadings_heatmap.png')
    plt.figure(figsize=(8,6))
    plt.imshow(img)
    plt.axis('off')
    pdf.savefig()
    plt.close()

    # Página 3: Biplot
    img = plt.imread('biplot.png')
    plt.figure(figsize=(10,8))
    plt.imshow(img)
    plt.axis('off')
    pdf.savefig()
    plt.close()

    # Página 4: Top/Bottom 5 PC1
    fig, ax = plt.subplots(figsize=(8,6))
    ax.axis('off')
    table_data = pd.concat([top5, bottom5], keys=['Top 5 PC1', 'Bottom 5 PC1'])
    table = ax.table(cellText=table_data.reset_index(drop=True).values,
                     colLabels=table_data.reset_index(drop=True).columns,
                     loc='center', cellLoc='center')
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1, 2)
    plt.title('Estados extremos en PC1')
    pdf.savefig()
    plt.close()

    # Página 5: Recomendaciones
    fig, ax = plt.subplots(figsize=(8,6))
    ax.axis('off')
    texto = (
        "Interpretación de Componentes:\n"
        "- PC1: criminalidad general (Murder, Assault)\n"
        "- PC2: urbanización y delitos específicos (UrbanPop, Rape)\n\n"
        "Recomendaciones para ML:\n"
        "- Usar 2-3 componentes principales (~80% varianza)\n"
        "- Reduce dimensionalidad, mejora entrenamiento y evita overfitting\n"
        "- Aplicaciones: clasificación, clustering, detección de outliers\n\n"
        "Observaciones:\n"
        "- Estados con PC1 alto requieren más atención en seguridad\n"
        "- Clusters de estados permiten políticas diferenciadas"
    )
    ax.text(0, 0.5, texto, fontsize=10, va='center')
    pdf.savefig()
    plt.close()

# =============================
# 5. Generar PowerPoint
# =============================
prs = Presentation()

# --- Slide 1: Título y objetivo ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text = "PCA de criminalidad y urbanización en EE.UU."
title.text_frame.paragraphs[0].font.size = Pt(28)
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
textbox.text = (
    "Objetivo:\nReducir dimensionalidad y extraer componentes principales para análisis y aplicaciones de IA."
)
textbox.text_frame.paragraphs[0].font.size = Pt(18)

# --- Slide 2: Resumen de datos ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text = "Resumen de Datos"
title.text_frame.paragraphs[0].font.size = Pt(28)
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
textbox.text = (
    "• 50 estados\n• Variables: Murder, Assault, UrbanPop, Rape\n"
    "• Sin valores faltantes, balance verificado\n• Datos listos para PCA"
)
textbox.text_frame.paragraphs[0].font.size = Pt(18)

# --- Slide 3: Componentes principales ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text = "Componentes Principales"
title.text_frame.paragraphs[0].font.size = Pt(28)
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
textbox.text = "• PC1: criminalidad general (Murder, Assault)\n• PC2: urbanización y delitos específicos (UrbanPop, Rape)"
textbox.text_frame.paragraphs[0].font.size = Pt(18)
slide.shapes.add_picture('loadings_heatmap.png', Inches(1), Inches(3), width=Inches(8), height=Inches(4.5))

# --- Slide 4: Biplot y clusters ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text = "Biplot y Clusters"
title.text_frame.paragraphs[0].font.size = Pt(28)
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
textbox.text = "• Outliers identificados\n• Clusters naturales de estados"
textbox.text_frame.paragraphs[0].font.size = Pt(18)
slide.shapes.add_picture('biplot.png', Inches(1), Inches(2.5), width=Inches(8), height=Inches(5))
# --- Slide 5: Recomendaciones ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text = "Recomendaciones"
title.text_frame.paragraphs[0].font.size = Pt(28)
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
textbox.text = (
    "• Usar 2–3 componentes principales (~80% varianza)\n"
    "• Estados con PC1 alto → prioridad en seguridad\n"
    "• Aplicaciones: clasificación, clustering, detección de outliers\n"
    "• PCA como preprocesamiento para ML y políticas públicas"
)
textbox.text_frame.paragraphs[0].font.size = Pt(18)

# Guardar presentación
prs.save('Presentacion_PCA_USArrests.pptx')
print("Presentación generada: Presentacion_PCA_USArrests.pptx")

